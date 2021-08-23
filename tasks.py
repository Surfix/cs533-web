import os
import sys
from pathlib import Path
import tarfile
import shutil
import re

from invoke import task

BRANCH = 'archive'
CURRENT = 'F21'
BUILD_DIR = Path('_build')
FULL_DIR = BUILD_DIR / 'site'
CUR_DIR = BUILD_DIR / 'dirhtml'
CUR_DEPLOY_DIR = FULL_DIR / CURRENT


def _msg(fmt, *args):
    print(fmt % args, file=sys.stderr)


@task
def build(c, rebuild=False, builder='dirhtml'):
    """
    Build the current site.
    """
    if rebuild:
        c.run(f'jb build --all --builder {builder} .')
    else:
        c.run(f'jb build --builder {builder} .')


@task
def serve(c, full_site=False, port=8000, rebuild=False):
    from livereload import Server

    def do_rebuild():
        _msg('building and copying site')
        build(c, rebuild=rebuild)
        if full_site:
            build_site(c)

    _ignore_re = re.compile(r'^(\.[\\/])?_build[\\/]')

    def ignore(p):
        return _ignore_re.match(p)

    do_rebuild()
    srv_root = FULL_DIR if full_site else CUR_DIR
    server = Server()
    server.watch('.', do_rebuild, ignore=ignore)
    server.serve(port=port, root=srv_root)


@task
def copy_archive(c):
    archive_path = BUILD_DIR / 'archive.tar'
    p = archive_path.as_posix()

    FULL_DIR.mkdir(exist_ok=True, parents=True)

    _msg('creating git archive')
    c.run(f'git archive -o {p} archive')

    with tarfile.open(archive_path) as tf:
        _msg('extracting %s', archive_path)
        tf.extractall(FULL_DIR)


@task(pre=[copy_archive])
def copy_root(c):
    _msg('copying static root files')
    shutil.copytree('_root', FULL_DIR, dirs_exist_ok=True)


@task(pre=[build, copy_root])
def build_site(c):
    _msg('copying current tree to site target')
    shutil.copytree(CUR_DIR, CUR_DEPLOY_DIR, dirs_exist_ok=True)


@task
def clean(c):
    _msg('removing build directory')
    shutil.rmtree(BUILD_DIR)


@task
def extract_slides(c, dir):
    out_dir = Path('videos')
    dir = Path(dir)
    import pptx
    for file in dir.glob("*.pptx"):
        print('reading', file.name)
        name = file.stem
        out_file = out_dir / f'{name}.slides.txt'
        ppt = pptx.Presentation(file)
        with open(out_file, 'w', encoding='utf8') as otf:
            for slide in ppt.slides:
                for shape in slide.shapes:
                    text = getattr(shape, 'text', None)
                    if text:
                        print(text, file=otf)
